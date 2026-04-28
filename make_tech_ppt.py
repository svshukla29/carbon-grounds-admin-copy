# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
ACCENT    = RGBColor(0x2E, 0xCC, 0x71)   # carbon green
ACCENT2   = RGBColor(0x1A, 0xBC, 0x9C)   # teal
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xBB, 0xCC)
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)
RED_LIGHT = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


def add_bg(slide, color=BG_DARK):
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, bg=None, text="", fontsize=18,
        bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg:
        from pptx.util import Pt as Pt2
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox


def rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def title_slide(text, sub="", tag=""):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    # Green accent bar left
    rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    # Green bar bottom
    rect(slide, 0, 6.8, 13.33, 0.08, ACCENT)
    box(slide, 0.3, 0.3, 4, 0.5, text="🌱 CARBON GROUNDS", fontsize=13, color=ACCENT, bold=True)
    box(slide, 0.3, 1.5, 12, 1.5, text=text, fontsize=40, bold=True, color=WHITE)
    if sub:
        box(slide, 0.3, 3.2, 12, 0.8, text=sub, fontsize=20, color=GRAY)
    if tag:
        box(slide, 0.3, 4.2, 12, 0.5, text=tag, fontsize=15, color=ACCENT, italic=True)
    box(slide, 0.3, 6.9, 12, 0.4, text="Vivek Shukla  |  April 2026", fontsize=12, color=GRAY)


def section_slide(num, title, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide)
    rect(slide, 0, 0, 13.33, 0.08, ACCENT)
    rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    box(slide, 0.3, 0.15, 3, 0.4, text=f"SLIDE {num}", fontsize=11, color=ACCENT, bold=True)
    box(slide, 0.3, 0.7, 12, 1.0, text=title, fontsize=32, bold=True, color=WHITE)
    if subtitle:
        box(slide, 0.3, 1.85, 12, 0.5, text=subtitle, fontsize=16, color=GRAY, italic=True)
    return slide


def bullet_slide(num, title, bullets, subtitle=""):
    slide = section_slide(num, title, subtitle)
    y = 2.5
    for b in bullets:
        if b.startswith("##"):
            box(slide, 0.5, y, 12, 0.4, text=b[2:].strip(), fontsize=14, color=ACCENT, bold=True)
            y += 0.45
        else:
            box(slide, 0.5, y, 12, 0.42, text="▸  " + b, fontsize=15, color=WHITE)
            y += 0.44
    return slide


def two_col(num, title, left_items, right_items, left_head="", right_head=""):
    slide = section_slide(num, title)
    if left_head:
        box(slide, 0.5, 2.4, 5.8, 0.4, text=left_head, fontsize=14, color=ACCENT, bold=True)
    if right_head:
        box(slide, 7.0, 2.4, 5.8, 0.4, text=right_head, fontsize=14, color=ACCENT2, bold=True)
    rect(slide, 6.6, 2.3, 0.05, 4.8, GRAY)
    y = 2.9
    for item in left_items:
        box(slide, 0.5, y, 5.8, 0.4, text="▸  " + item, fontsize=14, color=WHITE)
        y += 0.43
    y = 2.9
    for item in right_items:
        box(slide, 7.0, y, 5.8, 0.4, text="▸  " + item, fontsize=14, color=WHITE)
        y += 0.43
    return slide


def table_slide(num, title, headers, rows, subtitle=""):
    slide = section_slide(num, title, subtitle)
    col_w = 12.5 / len(headers)
    x0 = 0.4
    y0 = 2.45
    # header row
    for i, h in enumerate(headers):
        rect(slide, x0 + i*col_w, y0, col_w - 0.05, 0.42, ACCENT)
        box(slide, x0 + i*col_w + 0.1, y0 + 0.04, col_w - 0.2, 0.35,
            text=h, fontsize=13, bold=True, color=BG_DARK)
    for r, row in enumerate(rows):
        row_bg = RGBColor(0x14, 0x2A, 0x3F) if r % 2 == 0 else RGBColor(0x0D, 0x1B, 0x2A)
        for c, cell in enumerate(row):
            rect(slide, x0 + c*col_w, y0 + 0.45 + r*0.42, col_w - 0.05, 0.40, row_bg)
            box(slide, x0 + c*col_w + 0.1, y0 + 0.48 + r*0.42, col_w - 0.2, 0.35,
                text=cell, fontsize=12, color=WHITE)
    return slide


# ══════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════

# 1 – Title
title_slide(
    "Web & Mobile App\nDevelopment Proposal",
    sub="Technical Architecture · Backend · Frontend · Mobile App",
    tag="Timeline: 1 Month  |  Platform: Web + iOS/Android"
)

# 2 – Agenda
bullet_slide(2, "Agenda", [
    "Project Vision & Purpose",
    "What We Have Already Built",
    "Current Tech Stack & Architecture",
    "Admin Web App — Completed Modules",
    "Backend REST API — Endpoints Overview",
    "Live Deployment on AWS",
    "Identified Gaps & Improvement Areas",
    "Proposed Backend Enhancements",
    "Proposed Frontend Enhancements",
    "Mobile App Proposal (React Native)",
    "1-Month Development Timeline",
    "Security Improvements",
    "Success Metrics & KPIs",
    "Risk Assessment",
    "Next Steps",
])

# 3 – Project Vision
slide = section_slide(3, "Project Vision", "What problem are we solving?")
box(slide, 0.5, 2.4, 12, 0.5, text="Carbon Grounds is a carbon credit management platform connecting:", fontsize=16, color=GRAY)
items = [
    ("🌾  Farmers", "Practice sustainable/regenerative agriculture & earn credits", ACCENT),
    ("🤝  Partners", "NGOs, Govt bodies, Corporates funding green projects", ACCENT2),
    ("📊  Administrators", "Manage projects, track credits, generate reports", YELLOW),
]
y = 3.1
for icon, desc, col in items:
    rect(slide, 0.5, y, 3.5, 0.65, RGBColor(0x14, 0x2A, 0x3F))
    box(slide, 0.65, y+0.08, 3.2, 0.5, text=icon, fontsize=16, bold=True, color=col)
    box(slide, 4.2, y+0.1, 8.5, 0.5, text=desc, fontsize=15, color=WHITE)
    y += 0.85
box(slide, 0.5, 5.6, 12, 0.5,
    text="🎯  India's Net Zero 2070 commitment + Carbon markets growing at 25% CAGR globally",
    fontsize=14, color=ACCENT, italic=True)

# 4 – What We've Already Built
slide = section_slide(4, "What We Have Already Built", "Current completed deliverables")
items4 = [
    ("✅", "Admin Web Application", "Next.js 14 — fully functional dashboard with live API stats", ACCENT),
    ("✅", "Backend REST API", "NestJS — 30+ secured endpoints incl. file upload", ACCENT),
    ("✅", "Database", "PostgreSQL — 7 tables + new land/farmer columns", ACCENT),
    ("✅", "Authentication", "JWT login, refresh tokens, role-based access", ACCENT),
    ("✅", "Cloud Deployment", "Live on AWS EC2 — accessible from internet", ACCENT),
    ("✅", "Core Modules", "Farmers (Project ID, Type), Projects (Land Details), Reports (File Upload), Teams (Add User), Dashboard", ACCENT),
]
y = 2.4
for icon, head, desc, col in items4:
    rect(slide, 0.4, y, 12.5, 0.58, RGBColor(0x10, 0x24, 0x38))
    box(slide, 0.6, y+0.08, 0.5, 0.4, text=icon, fontsize=18, color=col)
    box(slide, 1.2, y+0.08, 3.5, 0.4, text=head, fontsize=15, bold=True, color=WHITE)
    box(slide, 4.8, y+0.08, 8.0, 0.4, text=desc, fontsize=14, color=GRAY)
    y += 0.68

# 5 – Tech Stack
table_slide(5, "Current Technology Stack", 
    ["Layer", "Technology", "Purpose"],
    [
        ["Frontend", "Next.js 14 + TypeScript", "Admin web app with SSR"],
        ["Styling", "Tailwind CSS + ShadCN UI", "Responsive UI components"],
        ["Backend", "NestJS (Node.js)", "REST API server"],
        ["Database", "PostgreSQL (AWS RDS)", "Relational data storage"],
        ["ORM", "TypeORM", "Database models & migrations"],
        ["Auth", "JWT + bcrypt + Passport.js", "Secure authentication"],
        ["Server", "AWS EC2 + Nginx + PM2", "Cloud hosting & proxy"],
        ["File Storage", "AWS S3 (planned)", "Images, PDFs"],
    ]
)

# 6 – Architecture
slide = section_slide(6, "System Architecture (Current)", "Request flow from browser to database")
arch_lines = [
    "  Browser / Admin User",
    "         ↓",
    "  AWS EC2  —  Nginx Reverse Proxy",
    "    ↙                  ↘",
    "Next.js App        NestJS REST API",
    "(Port 3000)           (Port 3001)",
    "                         ↓",
    "              AWS RDS  PostgreSQL",
    "           Database: carbon_grounds",
    "  Tables: users · farmers · projects",
    "          partners · reports · teams",
]
y = 2.3
for line in arch_lines:
    col = ACCENT if "↓" in line or "↙" in line or "↘" in line else (ACCENT2 if "AWS" in line or "Next" in line or "Nest" in line else WHITE)
    box(slide, 2.5, y, 9, 0.38, text=line, fontsize=14, color=col, align=PP_ALIGN.CENTER)
    y += 0.4

# 7 – Admin Web Modules
slide = section_slide(7, "Admin Web App — Completed Modules", "6 fully working modules in the dashboard")
modules = [
    ("🏠 Dashboard", "Live API stats: Active Farmers, Projects, Total Villages\nActivity feed + progress charts"),
    ("🌾 Farmers", "Add/Edit/Delete farmers\nProject ID column, Type (crops), Land & Plant Details tabs"),
    ("🏗️ Projects", "Project ID, Land Details (area/type/soil/water/elevation/GPS)\nPrimary farmer assignment, carbon credits"),
    ("📎 Reports + Upload", "Quarterly/Monthly/Verification/Audit\nFile attachment: PDF, Word, Excel, Images (20MB)"),
    ("👥 Teams + Add User", "Add team members & create new users\nRole: Admin, Project Manager, Analyst, Viewer"),
    ("📊 Live Dashboard", "Fetches real stats via dashboardApi.getStats()\nLoader state, no hardcoded values"),
]
positions = [(0.4, 2.4), (4.55, 2.4), (8.7, 2.4), (0.4, 5.0), (4.55, 5.0), (8.7, 5.0)]
for (left, top), (name, desc) in zip(positions, modules):
    rect(slide, left, top, 3.9, 2.2, RGBColor(0x10, 0x26, 0x40))
    box(slide, left+0.15, top+0.12, 3.6, 0.45, text=name, fontsize=15, bold=True, color=ACCENT)
    box(slide, left+0.15, top+0.65, 3.6, 1.4, text=desc, fontsize=12, color=WHITE)

# 8 – Backend Endpoints
table_slide(8, "Backend REST API — Endpoints",
    ["Method", "Endpoint", "Description"],
    [
        ["POST", "/api/auth/login", "Admin login → JWT tokens"],
        ["GET / POST", "/api/farmers", "List (with projectId) / Create farmer"],
        ["GET/PATCH/DEL", "/api/farmers/:id", "Get / Update / Delete farmer"],
        ["GET / POST", "/api/projects", "List / Create (with Land Details + farmerId)"],
        ["PATCH", "/api/projects/:id", "Update project + farmer assignment"],
        ["POST", "/api/projects/:id/farmers", "Bulk assign farmers to project"],
        ["GET / POST", "/api/reports", "List / Create report"],
        ["POST", "/api/reports/:id/upload", "Upload PDF/Word/Excel/Image (20MB)"],
        ["GET", "/api/reports/files/:filename", "Serve uploaded file"],
        ["GET", "/api/dashboard/stats", "activeFarmers + totalVillages + credits"],
    ],
    subtitle="All routes protected with JWT Bearer token"
)

# 9 – Live Deployment
slide = section_slide(9, "Live Deployment on AWS", "Already running in production")
info = [
    ("🌐 Frontend URL", "http://43.204.144.76", ACCENT),
    ("⚙️  API Base URL", "http://43.204.144.76/api", ACCENT2),
    ("🗄️  Database", "AWS RDS PostgreSQL  |  ap-south-1", YELLOW),
    ("🖥️  Server", "AWS EC2 (Nginx + PM2)", WHITE),
]
y = 2.5
for label, value, col in info:
    rect(slide, 0.5, y, 12.3, 0.65, RGBColor(0x10, 0x26, 0x40))
    box(slide, 0.8, y+0.1, 3.5, 0.44, text=label, fontsize=15, bold=True, color=col)
    box(slide, 4.5, y+0.1, 8.0, 0.44, text=value, fontsize=15, color=WHITE)
    y += 0.78
box(slide, 0.5, 5.5, 12.3, 0.5, text="✅  Login working  ·  ✅  All CRUD APIs functional  ·  ✅  Dashboard showing live data", fontsize=14, color=ACCENT, align=PP_ALIGN.CENTER)

# 10 – Gaps
table_slide(10, "Identified Gaps & Improvement Areas",
    ["Gap", "Impact", "Priority"],
    [
        ["File upload", "Report attachments (PDF/Word/Excel/Image)", "✅ Done"],
        ["Dashboard live stats", "activeFarmers + totalVillages from real API", "✅ Done"],
        ["Land Details on Projects", "area/type/soil/water/elevation/GPS columns", "✅ Done"],
        ["Add User from dashboard", "POST /users via Teams page dialog", "✅ Done"],
        ["No API pagination", "Performance issues at scale", "🔴 High"],
        ["No email notifications", "No alerts on status changes", "🟡 Medium"],
        ["No HTTPS / SSL", "Data not encrypted in transit", "🔴 High"],
        ["No mobile app", "Field officers cannot use system", "🔴 High"],
    ]
)

# 11 – Backend Enhancements
bullet_slide(11, "Proposed Backend Enhancements", [
    "##  Already Done (This Session)",
    "File upload endpoint: POST /reports/:id/upload  (multer, disk storage, 20MB limit)",
    "New DB columns: landArea, landType, soilType, waterSource, elevation, coordinates, farmerId on projects",
    "New DB column: projectId on farmers — synced when project farmerId is set",
    "Dashboard: activeFarmers (Verified count) + totalVillages (distinct locations)",
    "##  Week 1–2 (Next)",
    "Pagination & filtering on all GET endpoints  (page, limit, search, status)",
    "Email notifications module — Nodemailer + SMTP",
    "Carbon credit transaction history log per project",
    "##  Week 3–4",
    "Analytics API — monthly trends, regional farmer counts, project rates",
    "Audit log — track every create/update/delete with user + timestamp",
])

# 12 – Frontend Enhancements
bullet_slide(12, "Proposed Frontend Enhancements", [
    "##  Dashboard (Week 1–2)",
    "Line chart — Carbon credits earned over time",
    "India map — Project locations visualized geographically",
    "Pie chart — Project type distribution",
    "##  Data Management (Week 2)",
    "PDF export — Generate professional report PDFs",
    "CSV bulk import — Upload farmer/project data via spreadsheet",
    "Advanced search & date-range filters (server-side)",
    "##  UX & Security (Week 1 & 3)",
    "HTTPS / SSL — Let's Encrypt via Certbot + Nginx",
    "Settings page — Admin profile, change password",
    "Notification bell — In-app alert panel",
])

# 13 – Mobile App
slide = section_slide(13, "Mobile App Proposal", "React Native (Expo) — iOS + Android from one codebase")
targets = [
    ("📱 Field Officers", "Register farmers, capture photos,\nupdate status on-site", ACCENT),
    ("📊 Project Managers", "Monitor projects & credits\non the go", ACCENT2),
    ("🌾 Farmers (future)", "Check credit balance,\nproject updates", YELLOW),
]
x = 0.4
for name, desc, col in targets:
    rect(slide, x, 2.3, 3.9, 2.2, RGBColor(0x10, 0x26, 0x40))
    box(slide, x+0.15, 2.45, 3.6, 0.5, text=name, fontsize=16, bold=True, color=col)
    box(slide, x+0.15, 3.05, 3.6, 1.3, text=desc, fontsize=13, color=WHITE)
    x += 4.3
box(slide, 0.4, 4.9, 12.5, 0.55,
    text="🔑  Key advantage: Reuses the SAME NestJS backend — no new server development needed!",
    fontsize=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
screens = ["Login / Biometric Lock", "Dashboard summary cards", "Farmers list + Add form",
           "Camera photo capture", "GPS location auto-fill", "Project progress view",
           "Reports PDF viewer", "Push notifications"]
y = 5.6
x = 0.4
for i, s in enumerate(screens):
    box(slide, x, y, 3.0, 0.38, text="▸  " + s, fontsize=12, color=GRAY)
    x += 3.3
    if i == 3:
        x = 0.4; y = 6.05

# 14 – Mobile Architecture
slide = section_slide(14, "Mobile App — Technical Architecture")
lines = [
    ("React Native (Expo)  —  Single Codebase", WHITE, True),
    ("iOS App          Android App", GRAY, False),
    ("          ↕ Shared Code ↕", ACCENT, False),
    ("Navigation (React Navigation)", WHITE, False),
    ("State Management (Zustand + React Query)", WHITE, False),
    ("API calls via Axios  →  Same NestJS Backend", ACCENT2, False),
    ("Tokens stored in AsyncStorage (secure)", WHITE, False),
    ("Camera (Expo Camera)  ·  GPS (Expo Location)", WHITE, False),
    ("Offline storage (expo-sqlite)", YELLOW, False),
    ("          ↓", ACCENT, False),
    ("NestJS Backend on AWS  (unchanged)", ACCENT, True),
]
y = 2.3
for text, col, bold in lines:
    box(slide, 1.5, y, 10, 0.42, text=text, fontsize=15, bold=bold, color=col, align=PP_ALIGN.CENTER)
    y += 0.44

# 15 – Offline Mode
bullet_slide(15, "Mobile App — Offline Capability",
    [
        "##  Problem",
        "Field officers often work in remote areas with poor internet connectivity",
        "No internet means they cannot register farmers or update status in the field",
        "##  Solution: Offline-First Architecture",
        "All data cached locally using expo-sqlite on device",
        "Changes queued when offline — auto-sync when internet restores",
        "Conflict resolution: last-write-wins strategy with conflict log",
        "##  Libraries Used",
        "@tanstack/react-query with persistence layer",
        "expo-sqlite for local database on device",
        "NetInfo API to detect connectivity changes automatically",
    ],
    subtitle="Field officers often work in areas with poor connectivity"
)

# 16 – Timeline
slide = section_slide(16, "1-Month Development Timeline", "Week-by-week delivery plan")
weeks = [
    ("WEEK 1\nApr 28 – May 4", [
        "Pagination & filtering APIs",
        "AWS S3 file upload",
        "HTTPS / SSL setup",
        "Advanced dashboard charts",
    ], ACCENT),
    ("WEEK 2\nMay 5 – May 11", [
        "Email notifications module",
        "Carbon credit transaction log",
        "PDF report export",
        "Mobile: Auth + Navigation",
    ], ACCENT2),
    ("WEEK 3\nMay 12 – May 18", [
        "Analytics API",
        "Audit log module",
        "Mobile: Dashboard + Farmers",
        "Camera + GPS integration",
    ], YELLOW),
    ("WEEK 4\nMay 19 – May 25", [
        "Mobile: Reports + Notifications",
        "Offline sync implementation",
        "Full QA & testing",
        "Beta deployment (Expo Go)",
    ], RGBColor(0xAF, 0x7A, 0xC5)),
]
x = 0.3
for head, tasks, col in weeks:
    rect(slide, x, 2.3, 3.1, 4.8, RGBColor(0x10, 0x26, 0x40))
    box(slide, x+0.1, 2.4, 2.9, 0.8, text=head, fontsize=13, bold=True, color=col)
    y = 3.3
    for t in tasks:
        box(slide, x+0.1, y, 2.9, 0.5, text="▸  " + t, fontsize=12, color=WHITE)
        y += 0.55
    x += 3.28

# 17 – Security
table_slide(17, "Security Improvements",
    ["Feature", "Implementation", "Status"],
    [
        ["JWT Auth", "Access 15min + Refresh 7d rotation", "✅ Done"],
        ["Password hashing", "bcrypt with salt rounds", "✅ Done"],
        ["Role-based guards", "Admin / User roles on all routes", "✅ Done"],
        ["CORS protection", "Whitelist frontend origin only", "✅ Done"],
        ["HTTPS / SSL", "Let's Encrypt + Certbot + Nginx", "🔲 Week 1"],
        ["Rate limiting", "@nestjs/throttler — 100 req/min", "🔲 Week 1"],
        ["Helmet.js", "HTTP security headers", "🔲 Week 1"],
        ["Input validation", "class-validator on all DTOs", "🔲 Week 1"],
    ]
)

# 18 – Metrics
table_slide(18, "Success Metrics & KPIs",
    ["Metric", "Target", "When"],
    [
        ["API response time (p95)", "< 200ms", "End of Month 1"],
        ["API test coverage", "> 80%", "End of Month 1"],
        ["Mobile app crash rate", "< 1%", "Beta release"],
        ["System uptime", "> 99%", "Ongoing"],
        ["Farmers onboarded", "100+", "3-month view"],
        ["Projects created", "20+", "3-month view"],
        ["Carbon credits tracked", "10,000+", "3-month view"],
        ["Mobile installs (field officers)", "50+", "3-month view"],
    ]
)

# 19 – Risks
table_slide(19, "Risk Assessment & Mitigation",
    ["Risk", "Likelihood", "Mitigation"],
    [
        ["AWS cost overrun", "Medium", "Set billing alerts, use free tier"],
        ["Mobile store rejection", "Low", "Follow guidelines, test early"],
        ["DB performance at scale", "Low", "Add indexes, query optimization"],
        ["Offline sync conflicts", "Medium", "Last-write-wins + conflict log"],
        ["Poor field connectivity", "High", "Full offline mode (SQLite)"],
        ["Security breach", "Low", "Rate limiting, SSL, audit logs"],
    ]
)

# 20 – Next Steps
slide = section_slide(20, "Next Steps & Ask", "What happens after this meeting?")
steps = [
    ("✅  1", "Approve this proposal", ACCENT),
    ("📅  2", "Schedule weekly Friday check-ins", ACCENT2),
    ("🔑  3", "Provide: Domain name, Apple/Google dev accounts", YELLOW),
    ("🚀  4", "Development begins Week 1 immediately", WHITE),
]
y = 2.5
for num, text, col in steps:
    rect(slide, 0.4, y, 12.5, 0.72, RGBColor(0x10, 0x26, 0x40))
    box(slide, 0.7, y+0.12, 0.8, 0.48, text=num, fontsize=16, bold=True, color=col)
    box(slide, 1.8, y+0.12, 10.8, 0.48, text=text, fontsize=17, color=WHITE)
    y += 0.85
box(slide, 0.4, 6.1, 12.5, 0.55,
    text='"I will deliver a fully functional Web + Mobile platform within 4 weeks."',
    fontsize=16, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)

# Save
out = r"c:\Users\svshu\Desktop\CarbonGrounds_Technical_Proposal.pptx"
prs.save(out)
print("DONE: " + out)
