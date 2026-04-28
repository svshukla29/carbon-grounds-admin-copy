# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG   = RGBColor(0x0D,0x1B,0x2A)
GRN  = RGBColor(0x2E,0xCC,0x71)
TEAL = RGBColor(0x1A,0xBC,0x9C)
WHT  = RGBColor(0xFF,0xFF,0xFF)
GRY  = RGBColor(0xAA,0xBB,0xCC)
YLW  = RGBColor(0xF3,0x9C,0x12)
BLU  = RGBColor(0x5D,0xAD,0xE8)
DRK  = RGBColor(0x10,0x24,0x38)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLK = prs.slide_layouts[6]

def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def bx(slide, l, t, w, h, txt="", fs=16, bold=False, col=WHT,
        align=PP_ALIGN.LEFT, italic=False, bg_col=None):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = txt
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col
    if bg_col:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg_col
    return tb

def rx(slide, l, t, w, h, col, line=False):
    s = slide.shapes.add_shape(1,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = col
    if line: s.line.color.rgb = GRN
    else: s.line.fill.background()
    return s

def hdr(slide, num, title, sub=""):
    bg(slide)
    rx(slide,0,0,13.33,0.07,GRN); rx(slide,0,0,0.07,7.5,GRN)
    bx(slide,0.3,0.12,3,0.35,f"SLIDE {num}",11,True,GRN)
    bx(slide,0.3,0.65,12.5,1.0,title,30,True,WHT)
    if sub: bx(slide,0.3,1.8,12.5,0.45,sub,15,False,GRY,italic=True)
    return slide

def bullets(slide, items, y0=2.4):
    y = y0
    for b in items:
        if b.startswith("##"):
            bx(slide,0.5,y,12,0.38,b[2:].strip(),13,True,GRN); y+=0.42
        else:
            bx(slide,0.5,y,12,0.4,"  "+b,14,False,WHT); y+=0.42

# ─── SLIDE 1: Title ────────────────────────────────────────
s = prs.slides.add_slide(BLK); bg(s)
rx(s,0,0,0.07,7.5,GRN); rx(s,0,6.8,13.33,0.07,GRN)
bx(s,0.3,0.3,6,0.5,"CARBON GROUNDS",14,True,GRN)
bx(s,0.3,1.2,12,1.8,"Digital Platform\nProposal",44,True,WHT)
bx(s,0.3,3.3,12,0.6,"Website + Mobile App — What We've Built & What's Coming Next",18,False,GRY)
bx(s,0.3,4.2,12,0.5,"Presented by: Vivek Shukla  |  April 2026",14,False,TEAL)
bx(s,0.3,6.85,12,0.4,"Confidential — For Internal Review",11,False,GRY,PP_ALIGN.RIGHT)

# ─── SLIDE 2: The Problem ──────────────────────────────────
s = hdr(prs.slides.add_slide(BLK),2,"The Problem We Are Solving","India's farmers are doing great work — but it's invisible")
probs = [
    ("Farmers doing sustainable farming get no recognition or financial reward","🌾"),
    ("Carbon projects have no central system to track progress or impact","📋"),
    ("Partners and donors have no visibility into how funds are being used","🤝"),
    ("All data is managed manually — slow, error-prone, hard to scale","📂"),
]
y=2.4
for txt,icon in probs:
    rx(s,0.4,y,12.5,0.65,DRK)
    bx(s,0.6,y+0.1,0.6,0.45,icon,20,False,GRN)
    bx(s,1.4,y+0.12,11.2,0.42,txt,15,False,WHT)
    y+=0.78

# ─── SLIDE 3: Our Solution ─────────────────────────────────
s = hdr(prs.slides.add_slide(BLK),3,"Our Solution — Carbon Grounds Platform","One digital system for everything")
bx(s,0.4,2.15,12.5,0.5,"A single web + mobile platform that connects farmers, partners, and administrators.",16,False,GRY)
cards = [
    ("Website\n(Admin Panel)","Manage all data from a computer\nAnywhere, anytime",GRN),
    ("Mobile App\n(Field Workers)","Register farmers on-site\nWork even without internet",TEAL),
    ("Live Dashboard","See all numbers in real-time\nFarmers, projects, villages",YLW),
]
x=0.4
for title,desc,col in cards:
    rx(s,x,2.9,3.9,3.5,DRK)
    bx(s,x+0.15,3.05,3.6,0.7,title,17,True,col)
    bx(s,x+0.15,3.85,3.6,2.0,desc,13,False,WHT)
    x+=4.3

# ─── SLIDE 4: Who Uses This? ───────────────────────────────
s = hdr(prs.slides.add_slide(BLK),4,"Who Uses the Platform?","Three types of users, one connected system")
users = [
    ("Administrators","Office staff who manage the entire platform\nAdd/edit farmers, projects, reports, users","Admin Panel (Website)","👨‍💼",GRN),
    ("Field Officers","On-ground staff who visit farmers\nRegister new farmers, update status on mobile","Mobile App","👷",TEAL),
    ("Partners & Donors","NGOs, government, corporates\nTrack how projects are progressing","View-only reports (future)","🤝",YLW),
]
y=2.3
for role,desc,tool,icon,col in users:
    rx(s,0.4,y,12.5,1.05,DRK)
    bx(s,0.6,y+0.15,0.5,0.75,icon,24,False,col)
    bx(s,1.3,y+0.12,4.0,0.4,role,16,True,col)
    bx(s,1.3,y+0.55,6.5,0.42,desc,12,False,WHT)
    rx(s,9.5,y+0.2,3.2,0.65,RGBColor(0x14,0x35,0x50))
    bx(s,9.6,y+0.28,3.0,0.42,tool,12,True,col)
    y+=1.2

# ─── SLIDE 5: Simple Workflow ──────────────────────────────
s = hdr(prs.slides.add_slide(BLK),5,"How It Works — Simple Workflow","From farmer registration to carbon credits")
steps = [
    ("1","Field officer visits farmer & registers them on mobile app",GRN),
    ("2","Admin creates a Carbon Project with Land Details and links farmers",TEAL),
    ("3","Project tracks land area, type, soil, water source, carbon absorbed",YLW),
    ("4","Reports are generated — Quarterly, Verification, Impact (with file attachments)",BLU),
    ("5","Admin can add users directly and manage team access",GRN),
    ("6","Carbon credits are tracked and can be traded/reported",TEAL),
]
y=2.35
for num,txt,col in steps:
    rx(s,0.4,y,0.55,0.58,col)
    bx(s,0.48,y+0.08,0.4,0.42,num,18,True,BG,PP_ALIGN.CENTER)
    rx(s,1.1,y,11.65,0.58,DRK)
    bx(s,1.25,y+0.1,11.3,0.42,txt,14,False,WHT)
    y+=0.73

# ─── SLIDE 6: What's Already Built ────────────────────────
s = hdr(prs.slides.add_slide(BLK),6,"What's Already Built & Working","Completed and live on the internet right now")
done = [
    ("Admin Login","Secure login page with username & password","✅"),
    ("Farmer Management","Add, edit, delete farmers — with location, type/crops, Project ID","✅"),
    ("Project Management","Create projects with full Land Details, link farmers, track credits","✅"),
    ("Reports with File Upload","Create reports & attach PDF/Word/Excel/Image files","✅"),
    ("Team Management","Add team members & create new users directly from dashboard","✅"),
    ("Live Dashboard","Real-time numbers — Active Farmers, Projects, Total Villages","✅"),
]
y=2.35
for feat,desc,icon in done:
    rx(s,0.4,y,12.5,0.62,DRK)
    bx(s,0.6,y+0.1,0.5,0.4,icon,16,True,GRN)
    bx(s,1.2,y+0.1,3.5,0.4,feat,14,True,WHT)
    bx(s,5.0,y+0.1,8.0,0.4,desc,13,False,GRY)
    y+=0.66

# ─── SLIDE 7: Live on Internet ────────────────────────────
s = hdr(prs.slides.add_slide(BLK),7,"Already Live on the Internet","The platform is deployed and running")
bx(s,0.4,2.1,12.5,0.5,"You can access it right now from any browser:",16,False,GRY)
rx(s,0.4,2.7,12.5,1.0,DRK)
bx(s,0.7,2.82,12,0.7,"http://43.204.144.76",22,True,GRN,PP_ALIGN.CENTER)
facts = [
    ("Hosted on","Amazon cloud servers — India region","🌐"),
    ("Database","Secure cloud database — all data safely stored","🗄️"),
    ("Security","Password protected — only authorized users can login","🔒"),
    ("Availability","Runs 24 hours, 7 days a week — always accessible","⏰"),
]
y=3.9
for label,desc,icon in facts:
    rx(s,0.4,y,12.5,0.65,DRK)
    bx(s,0.6,y+0.1,0.5,0.45,icon,20,False,TEAL)
    bx(s,1.3,y+0.1,3.0,0.4,label,14,True,WHT)
    bx(s,4.5,y+0.1,8.2,0.4,desc,13,False,GRY)
    y+=0.75

# ─── SLIDE 8: Farmer Module ────────────────────────────────
s = hdr(prs.slides.add_slide(BLK),8,"Farmer Management — What Can You Do?","Complete farmer profile and tracking")
items8=[
    "Add a new farmer with name, phone, location, land area",
    "Record what type of crops/plants they grow (shown as 'Type')",
    "Mark farmer as Verified (Active) or Pending verification",
    "Each farmer shows their Project ID for quick identification",
    "View full farmer profile — Land Details & Plant Details tabs",
    "Edit or delete farmer records anytime",
    "Search and filter farmers by status or location",
]
bullets(s,["##  Farmer Module Features"]+items8)

# ─── SLIDE 9: Project Module ───────────────────────────────
s = hdr(prs.slides.add_slide(BLK),9,"Project Management — Tracking Carbon Work","Each project tracks real carbon work on the ground")
ptypes = [("Agroforestry","Trees + crops grown together on farmland",GRN),
          ("Reforestation","Planting trees on degraded land",TEAL),
          ("Soil Carbon","Improving soil health to absorb more carbon",YLW),
          ("Blue Carbon","Mangroves, wetlands, coastal ecosystems",BLU)]
x=0.4
for pt,desc,col in ptypes:
    rx(s,x,2.3,3.0,1.8,DRK)
    bx(s,x+0.1,2.42,2.8,0.5,pt,14,True,col)
    bx(s,x+0.1,3.0,2.8,1.0,desc,12,False,WHT)
    x+=3.22
bullets(s,[
    "##  Each Project Now Records",
    "Project ID, location, type, carbon credits",
    "Land Area (ha), Land Type, Soil Type, Water Source, Elevation, GPS Coordinates",
    "Primary Farmer assigned to the project",
    "All reports linked to the project (with file attachments)",
])

# ─── SLIDE 10: Reports & Dashboard ───────────────────────
s = hdr(prs.slides.add_slide(BLK),10,"Reports & Dashboard","Data-driven decisions at a glance")
bx(s,0.4,2.1,6.0,0.4,"REPORTS",14,True,GRN)
rtypes=["Quarterly Summary — every 3 months",
        "Monthly Summary — every month",
        "Verification Report — third party check",
        "Impact Report — outcomes & results",
        "Audit Report — financial & compliance",
        "File attachments — PDF, Word, Excel, Images"]
y=2.6
for r in rtypes:
    bx(s,0.5,y,5.8,0.42,"  "+r,13,False,WHT); y+=0.44
rx(s,6.6,2.0,0.05,5.0,GRY)
bx(s,6.9,2.1,6.0,0.4,"DASHBOARD SHOWS",14,True,TEAL)
dstats=["Total active projects",
        "Active farmers (verified)",
        "Total villages covered",
        "Carbon credits generated",
        "Number of team members",
        "Recent activity feed"]
y=2.6
for d in dstats:
    bx(s,7.0,y,5.8,0.42,"  "+d,13,False,WHT); y+=0.44

# ─── SLIDE 11: What's Missing ────────────────────────────
s = hdr(prs.slides.add_slide(BLK),11,"What's Not There Yet — Our Next Steps","Areas identified for improvement in 1 month")
gaps=[
    ("Mobile App","Field officers need a phone app to work on-site","Weeks 2-4",GRN),
    ("Email Alerts","No automatic emails when a status changes","Week 2",TEAL),
    ("Better Charts","Need graphs for carbon trends over time","Week 1",BLU),
    ("Data Import","Currently adding farmers one by one — need bulk upload","Week 2",YLW),
    ("Secure Web Address","Need https:// instead of plain http://","Week 1",GRN),
    ("PDF Downloads","Generate professional PDF from any report","Week 2",TEAL),
]
y=2.35
for feat,desc,when,col in gaps:
    rx(s,0.4,y,12.5,0.62,DRK)
    bx(s,0.6,y+0.1,3.2,0.42,feat,14,True,col)
    bx(s,4.0,y+0.1,7.0,0.42,desc,13,False,WHT)
    rx(s,11.2,y+0.1,1.55,0.42,RGBColor(0x14,0x35,0x50))
    bx(s,11.25,y+0.12,1.45,0.38,when,11,True,col,PP_ALIGN.CENTER)
    y+=0.73

# ─── SLIDE 12: Mobile App Plan ───────────────────────────
s = hdr(prs.slides.add_slide(BLK),12,"Mobile App — For Field Workers","Register farmers directly on-site, even without internet")
bx(s,0.4,2.1,12.5,0.5,"A simple phone app (works on Android & iPhone) for field officers:",15,False,GRY)
mscreens=[
    ("Login Screen","Same secure login as website","📱"),
    ("Dashboard","Quick summary of all numbers","📊"),
    ("Add Farmer","Register new farmer with camera photo","📷"),
    ("GPS Location","Auto-detect farmer location on map","📍"),
    ("View Projects","Check project details and progress","🏗️"),
    ("Works Offline","Save data locally, sync when internet returns","📡"),
]
x=0.4; y=2.75
for i,(name,desc,icon) in enumerate(mscreens):
    rx(s,x,y,3.9,1.55,DRK)
    bx(s,x+0.15,y+0.1,0.5,0.45,icon,20,False,GRN)
    bx(s,x+0.75,y+0.12,3.0,0.4,name,14,True,WHT)
    bx(s,x+0.15,y+0.65,3.6,0.75,desc,12,False,GRY)
    x+=4.3
    if i==2: x=0.4; y=4.5

# ─── SLIDE 13: Offline Benefit ──────────────────────────
s = hdr(prs.slides.add_slide(BLK),13,"Key Benefit — Works Without Internet","Perfect for remote farm locations")
rx(s,0.4,2.2,12.5,1.3,DRK)
bx(s,0.7,2.35,12,0.5,"The Challenge: Many farms are in areas with no internet or weak signal",16,False,YLW)
bx(s,0.7,2.8,12,0.5,"Our Solution: The mobile app saves data on the phone first, then uploads automatically when internet is available",14,False,WHT)
steps13=[
    ("Field officer opens app","Works normally, no internet needed"),
    ("Registers farmer / updates status","Data saved on phone"),
    ("Returns to area with internet","App detects connection"),
    ("All data uploads automatically","Nothing is lost"),
]
y=3.75
for step,detail in steps13:
    rx(s,0.4,y,6.0,0.65,RGBColor(0x0D,0x2D,0x44))
    bx(s,0.6,y+0.1,5.7,0.45,step,14,True,GRN)
    rx(s,6.7,y,6.2,0.65,DRK)
    bx(s,6.9,y+0.1,5.8,0.45,detail,14,False,WHT)
    y+=0.82

# ─── SLIDE 14: 1 Month Plan ──────────────────────────────
s = hdr(prs.slides.add_slide(BLK),14,"1-Month Development Plan","Simple timeline of what gets built when")
wks=[
    ("WEEK 1","Apr 28 – May 4",[
        "Secure web address (HTTPS)",
        "Better dashboard charts",
        "Faster search & filters",
        "PDF report downloads",
    ],GRN),
    ("WEEK 2","May 5 – May 11",[
        "Email alerts & notifications",
        "Carbon credit history log",
        "Bulk farmer data import",
        "Analytics & trend charts",
    ],TEAL),
    ("WEEK 3","May 12 – May 18",[
        "Mobile app: Login & Dashboard",
        "Mobile app: Add farmer + photo",
        "Mobile app: View projects",
        "GPS location auto-fill",
    ],YLW),
    ("WEEK 4","May 19 – May 25",[
        "Mobile app: Offline mode",
        "Full testing on all devices",
        "Bug fixes & improvements",
        "Beta release for testing",
    ],BLU),
]
x=0.3
for wk,dates,tasks,col in wks:
    rx(s,x,2.25,3.1,5.0,DRK)
    rx(s,x,2.25,3.1,0.72,col)
    bx(s,x+0.1,2.3,2.9,0.38,wk,14,True,BG,PP_ALIGN.CENTER)
    bx(s,x+0.1,2.65,2.9,0.28,dates,10,False,BG,PP_ALIGN.CENTER)
    y2=3.1
    for t in tasks:
        bx(s,x+0.1,y2,2.9,0.5,"  "+t,12,False,WHT); y2+=0.54
    x+=3.27

# ─── SLIDE 15: What Client Gets ─────────────────────────
s = hdr(prs.slides.add_slide(BLK),15,"What You Will Have After 1 Month","Complete digital platform ready to scale")
deliverables=[
    ("Website (Admin Panel)","Fully working with all improvements — secure, fast, with charts & exports","✅"),
    ("Mobile App (Android + iPhone)","Field workers can register farmers, take photos, work offline","✅"),
    ("Email Notifications","Automatic alerts when farmer verified, project completed, etc.","✅"),
    ("PDF Reports","Download any report as a professional PDF document","✅"),
    ("Data Import Tool","Upload 100s of farmers at once from Excel/CSV file","✅"),
    ("Secure Website","https:// encrypted connection, data protected","✅"),
]
y=2.35
for feat,desc,icon in deliverables:
    rx(s,0.4,y,12.5,0.62,DRK)
    bx(s,0.6,y+0.1,0.5,0.42,icon,16,True,GRN)
    bx(s,1.2,y+0.1,3.5,0.42,feat,14,True,WHT)
    bx(s,5.0,y+0.1,7.9,0.42,desc,13,False,GRY)
    y+=0.72

# ─── SLIDE 16: Results Expected ─────────────────────────
s = hdr(prs.slides.add_slide(BLK),16,"What Results Can We Expect?","Goals for the first 3 months after launch")
results=[
    ("100+ Farmers","Registered and verified on the platform","🌾"),
    ("20+ Projects","Carbon projects with full land details tracked","🏗️"),
    ("50+ Field Officers","Using the mobile app on Android/iPhone","📱"),
    ("Automated Reports","No more manual report preparation","📊"),
    ("Partner Visibility","Partners can track fund utilization","🤝"),
    ("Scalable Platform","Ready to grow to 1000s of farmers","📈"),
]
x=0.4; y=2.35
for i,(metric,desc,icon) in enumerate(results):
    rx(s,x,y,3.9,1.75,DRK)
    bx(s,x+0.15,y+0.15,0.55,0.5,icon,22,False,GRN)
    bx(s,x+0.85,y+0.17,2.9,0.45,metric,15,True,WHT)
    bx(s,x+0.15,y+0.75,3.6,0.85,desc,12,False,GRY)
    x+=4.3
    if i==2: x=0.4; y=4.3

# ─── SLIDE 17: Next Steps ────────────────────────────────
s = hdr(prs.slides.add_slide(BLK),17,"Next Steps — What We Need to Start","Simple action items to kick off development")
actions=[
    ("1","Approve this proposal so development can begin","Your approval",GRN),
    ("2","Share a preferred domain name for the website","e.g. carbongrounds.in",TEAL),
    ("3","Decide if mobile app needs Apple App Store or just Android","Affects timeline by 1 week",YLW),
    ("4","Provide list of initial farmers/projects for data entry","For testing purposes",BLU),
    ("5","Schedule weekly 30-min update call every Friday","To track progress",GRN),
]
y=2.35
for num,action,note,col in actions:
    rx(s,0.4,y,12.5,0.72,DRK)
    rx(s,0.4,y,0.55,0.72,col)
    bx(s,0.47,y+0.14,0.42,0.45,num,18,True,BG,PP_ALIGN.CENTER)
    bx(s,1.15,y+0.14,8.0,0.45,action,15,False,WHT)
    bx(s,9.4,y+0.14,3.4,0.45,note,12,False,GRY,PP_ALIGN.RIGHT)
    y+=0.86

# ─── SLIDE 18: Thank You ─────────────────────────────────
s = prs.slides.add_slide(BLK); bg(s)
rx(s,0,0,0.07,7.5,GRN); rx(s,0,6.8,13.33,0.07,GRN)
rx(s,1.5,1.5,10.33,4.5,DRK)
bx(s,1.7,1.8,10,0.6,"Thank You",38,True,GRN,PP_ALIGN.CENTER)
bx(s,1.7,2.7,10,0.5,"Vivek Shukla  —  Full Stack Developer",18,False,WHT,PP_ALIGN.CENTER)
bx(s,1.7,3.35,10,0.9,"Ready to build Carbon Grounds into a world-class\ndigital platform for sustainable agriculture.",16,False,GRY,PP_ALIGN.CENTER)
bx(s,1.7,4.4,10,0.45,"Questions & Discussion Welcome",14,True,TEAL,PP_ALIGN.CENTER)
bx(s,0.3,6.85,12.7,0.4,"Carbon Grounds  |  April 2026  |  Confidential",11,False,GRY,PP_ALIGN.CENTER)

out = r"c:\Users\svshu\Desktop\CarbonGrounds_Client_Proposal.pptx"
prs.save(out)
print("DONE: " + out)
